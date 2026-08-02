"""Generate a visual PawPal+ presentation deck (PawPal_Plus_Presentation.pptx)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---- Palette (paw-care theme) -------------------------------------------------
NAVY = RGBColor(0x14, 0x2A, 0x3F)      # deep slate — primary text / title bg
TEAL = RGBColor(0x14, 0x8F, 0x93)      # accent 1
ORANGE = RGBColor(0xF2, 0x8C, 0x28)    # accent 2
LIGHT = RGBColor(0xF4, 0xF7, 0xF8)     # slide background
MIST = RGBColor(0xE3, 0xEC, 0xEF)      # card background
INK = RGBColor(0x23, 0x33, 0x3D)       # body text
MUTE = RGBColor(0x5C, 0x6B, 0x73)      # secondary text
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def rect(slide, x, y, w, h, color, line=None, shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, x, y, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp


def text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=6, line_spacing=1.0):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, (txt, size, color, bold, *rest) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.line_spacing = line_spacing
        if rest:
            p.level = rest[0]
        r = p.add_run()
        r.text = txt
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = "Calibri"
    return tb


def header(slide, kicker, title, n):
    """Standard content-slide header with accent bar + page number."""
    rect(slide, 0, 0, SW, Inches(1.35), WHITE)
    rect(slide, Inches(0.6), Inches(0.42), Inches(0.14), Inches(0.62), ORANGE)
    text(slide, Inches(0.9), Inches(0.32), Inches(10.5), Inches(0.35),
         [(kicker.upper(), 12, TEAL, True)], space_after=0)
    text(slide, Inches(0.9), Inches(0.6), Inches(11.0), Inches(0.6),
         [(title, 28, NAVY, True)], space_after=0)
    # page marker
    dot = rect(slide, SW - Inches(1.05), Inches(0.45), Inches(0.55), Inches(0.55),
               TEAL, shape=MSO_SHAPE.OVAL)
    tf = dot.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = f"{n:02d}"
    r.font.size = Pt(16); r.font.bold = True; r.font.color.rgb = WHITE
    rect(slide, 0, SH - Inches(0.35), SW, Inches(0.35), NAVY)
    text(slide, Inches(0.6), SH - Inches(0.34), Inches(8), Inches(0.3),
         [("PawPal+  ·  Applied AI Pet-Care Planner", 10, WHITE, False)], space_after=0)


def bullets(slide, items, x=Inches(0.9), y=Inches(1.7), w=Inches(11.5), h=Inches(5.2),
            size=18, gap=10):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        lvl = 0
        if isinstance(it, tuple):
            it, lvl = it
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap); p.line_spacing = 1.05; p.level = lvl
        bullet = "▸  " if lvl == 0 else "–  "
        rb = p.add_run(); rb.text = bullet
        rb.font.size = Pt(size); rb.font.bold = True
        rb.font.color.rgb = ORANGE if lvl == 0 else TEAL
        rt = p.add_run(); rt.text = it
        rt.font.size = Pt(size if lvl == 0 else size - 2)
        rt.font.color.rgb = INK if lvl == 0 else MUTE
        rt.font.bold = False
    return tb


def card(slide, x, y, w, h, title, body, accent):
    rect(slide, x, y, w, h, MIST)
    rect(slide, x, y, w, Inches(0.09), accent)
    text(slide, x + Inches(0.25), y + Inches(0.22), w - Inches(0.5), Inches(0.5),
         [(title, 15, NAVY, True)], space_after=2)
    text(slide, x + Inches(0.25), y + Inches(0.72), w - Inches(0.5), h - Inches(0.9),
         [(body, 12.5, MUTE, False)], space_after=0, line_spacing=1.05)


# ============================================================ SLIDE 1 — TITLE
s = prs.slides.add_slide(BLANK)
bg(s, NAVY)
rect(s, 0, 0, SW, SH, NAVY)
# accent band
rect(s, 0, SH - Inches(2.6), SW, Inches(0.14), TEAL)
rect(s, 0, SH - Inches(2.46), SW, Inches(0.07), ORANGE)
# paw glyph
paw = rect(s, Inches(0.85), Inches(1.5), Inches(1.2), Inches(1.2), TEAL, shape=MSO_SHAPE.OVAL)
pf = paw.text_frame; pf.paragraphs[0].alignment = PP_ALIGN.CENTER
pr = pf.paragraphs[0].add_run(); pr.text = "🐾"; pr.font.size = Pt(44)
text(s, Inches(0.85), Inches(3.05), Inches(11.6), Inches(1.6),
     [("PawPal+", 60, WHITE, True)], space_after=0)
text(s, Inches(0.9), Inches(4.25), Inches(11.6), Inches(0.9),
     [("An Applied AI Pet-Care Planner with knowledge-based prioritization,", 22, MIST, False),
      ("traceable reasoning, and confidence scoring", 22, MIST, False)],
     space_after=2, line_spacing=1.1)
text(s, Inches(0.9), SH - Inches(2.15), Inches(11.6), Inches(0.6),
     [("Extended from the Module 2 pet-scheduling starter  ·  Built by Giridhar", 15, TEAL, True)],
     space_after=0)

# ============================================================ SLIDE 2 — PROBLEM
s = prs.slides.add_slide(BLANK); bg(s, LIGHT)
header(s, "Why it matters", "The Problem", 2)
bullets(s, [
    "Pet owners juggle many daily care tasks across multiple pets.",
    "Safety-critical routines — medication, feeding, exercise — are easy to miss.",
    "Most tools give a schedule but never explain WHY a plan was chosen…",
    "…or HOW confident the system is that the plan is right.",
], y=Inches(1.9), size=20, gap=16)
card(s, Inches(8.6), Inches(1.9), Inches(4.1), Inches(4.4),
     "The gap", "A scheduler that is not just automatic, but transparent and "
     "trustworthy — one that prioritizes what matters for a pet's health and "
     "shows its reasoning.", ORANGE)

# ============================================================ SLIDE 3 — SOLUTION
s = prs.slides.add_slide(BLANK); bg(s, LIGHT)
header(s, "What we built", "The Solution", 3)
bullets(s, [
    "A Streamlit app that collects owner, pet, and task details.",
    "Builds a daily plan from priority, preferred time, and available minutes.",
    "An AI planner guides scheduling using pet-care category knowledge.",
    "Every decision is logged in a reasoning trace with a confidence score.",
], y=Inches(1.9), w=Inches(7.4), size=19, gap=15)
card(s, Inches(8.6), Inches(1.9), Inches(4.1), Inches(2.1),
     "Input", "Owner availability, pets, and care tasks (title, duration, "
     "priority, category, preferred time).", TEAL)
card(s, Inches(8.6), Inches(4.2), Inches(4.1), Inches(2.1),
     "Output", "An ordered daily plan + skipped tasks, conflict warnings, "
     "confidence score, and full planner trace.", ORANGE)

# ============================================================ SLIDE 4 — ARCH
s = prs.slides.add_slide(BLANK); bg(s, LIGHT)
header(s, "How it fits together", "System Architecture", 4)


def node(x, y, w, label, sub, color, txtcolor=WHITE):
    box = rect(s, x, y, w, Inches(0.95), color, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label; r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = txtcolor
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = sub; r2.font.size = Pt(10); r2.font.color.rgb = txtcolor
    return box


def arrow(x, y, w):
    a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, w, Inches(0.35))
    a.fill.solid(); a.fill.fore_color.rgb = MUTE; a.line.fill.background(); a.shadow.inherit = False


row_y = Inches(2.3)
node(Inches(0.7), row_y, Inches(2.3), "User / Owner", "enters pets & tasks", NAVY)
arrow(Inches(3.05), row_y + Inches(0.3), Inches(0.55))
node(Inches(3.65), row_y, Inches(2.3), "Streamlit UI", "app.py", TEAL)
arrow(Inches(6.0), row_y + Inches(0.3), Inches(0.55))
node(Inches(6.6), row_y, Inches(2.5), "PetCarePlanner", "boosts · trace · confidence", ORANGE)
arrow(Inches(9.2), row_y + Inches(0.3), Inches(0.55))
node(Inches(9.8), row_y, Inches(2.7), "Scheduler", "ordering · slot search", NAVY)

# second row — supporting components
row2 = Inches(4.1)
node(Inches(6.6), row2, Inches(2.5), "CareKnowledgeBase", "category guidance (RAG-style)", TEAL)
node(Inches(9.8), row2, Inches(2.7), "DailyPlan", "scheduled + skipped tasks", NAVY)
node(Inches(3.65), row2, Inches(2.3), "Automated Tests", "reliability checks", ORANGE)
# connectors (thin lines via narrow rectangles)
rect(s, Inches(7.7), row_y + Inches(0.95), Inches(0.04), Inches(0.65), MUTE)   # planner->kb
rect(s, Inches(11.0), row_y + Inches(0.95), Inches(0.04), Inches(0.65), MUTE)  # scheduler->plan
text(s, Inches(0.7), Inches(5.4), Inches(11.9), Inches(1.2),
     [("Data flow:  input → knowledge retrieval → priority boost → plan build → "
       "conflict check → confidence.  Tests verify planner and scheduler behavior.",
       14, MUTE, False)], space_after=0, line_spacing=1.1)
text(s, Inches(0.7), Inches(6.15), Inches(11.9), Inches(0.5),
     [("Source: diagrams/system_architecture.mmd", 11, TEAL, True)], space_after=0)

# ============================================================ SLIDE 5 — AI FEATURES
s = prs.slides.add_slide(BLANK); bg(s, LIGHT)
header(s, "The intelligence", "AI Features", 5)
cw, ch = Inches(5.9), Inches(2.15)
card(s, Inches(0.7), Inches(1.85), cw, ch, "Retrieval (RAG-style)",
     "CareKnowledgeBase returns category guidance for each task before the plan "
     "is built — retrieval actively drives the priority boost.", TEAL)
card(s, Inches(6.9), Inches(1.85), cw, ch, "Knowledge-based boosts",
     "Medication (1.0) › Feeding (0.7) › Walk (0.5) › Grooming (0.3) reorder "
     "tasks so safety-critical care comes first.", ORANGE)
card(s, Inches(0.7), Inches(4.2), cw, ch, "Agentic workflow",
     "Plan → act → check → resolve: detects same-start conflicts and runs a "
     "resolution pass that re-prioritizes medication.", NAVY)
card(s, Inches(6.9), Inches(4.2), cw, ch, "Trace + confidence",
     "Every step is logged to an auditable trace, and a 0–1 confidence score "
     "reflects schedule coverage and remaining conflicts.", TEAL)

# ============================================================ SLIDE 6 — RELIABILITY
s = prs.slides.add_slide(BLANK); bg(s, LIGHT)
header(s, "Proving it works", "Reliability & Testing", 6)
# big stat tiles
def stat(x, num, label, color):
    rect(s, x, Inches(1.9), Inches(2.6), Inches(1.9), color)
    text(s, x, Inches(2.1), Inches(2.6), Inches(1.0),
         [(num, 44, WHITE, True)], align=PP_ALIGN.CENTER, space_after=0)
    text(s, x, Inches(3.15), Inches(2.6), Inches(0.6),
         [(label, 13, WHITE, False)], align=PP_ALIGN.CENTER, space_after=0)
stat(Inches(0.7), "11/11", "automated tests pass", TEAL)
stat(Inches(3.5), "1.00", "confidence on feasible plans", ORANGE)
stat(Inches(6.3), "4", "human-eval cases reviewed", NAVY)
bullets(s, [
    "pytest covers planner boosts, scheduling feasibility, conflict detection & skips.",
    "Confidence drops automatically when a task can't fit the available time.",
    "Guardrails: empty pet name and no-task cases are handled without crashing.",
    "Human-eval table documented in README (parseable markdown).",
], y=Inches(4.15), size=16, gap=11)

# ============================================================ SLIDE 7 — DEMO
s = prs.slides.add_slide(BLANK); bg(s, LIGHT)
header(s, "See it run", "Demo — Sample Output", 7)
# terminal-style card
rect(s, Inches(0.7), Inches(1.85), Inches(11.9), Inches(4.6), NAVY,
     shape=MSO_SHAPE.ROUNDED_RECTANGLE)
rect(s, Inches(0.7), Inches(1.85), Inches(11.9), Inches(0.5), RGBColor(0x0C, 0x1B, 0x28),
     shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, Inches(1.0), Inches(1.9), Inches(6), Inches(0.4),
     [("$ python main.py", 13, TEAL, True)], space_after=0)
demo = (
    "Today's Schedule\n"
    "====================\n"
    "08:00 — Medication (15 min, high priority)\n"
    "08:15 — Morning walk (30 min, high priority)\n"
    "08:45 — Brushing (12 min, medium priority)\n"
    "09:00 — Feeding (10 min, high priority)\n\n"
    "Planner confidence: 1.00\n\n"
    "- Retrieved guidance for 'Medication' (medication)…\n"
    "- Applied knowledge boost of 1.0 to 'Medication'.\n"
    "- Built initial plan: 4 scheduled, 0 skipped.\n"
    "- No conflicts detected.  Final confidence: 1.00"
)
tb = s.shapes.add_textbox(Inches(1.0), Inches(2.45), Inches(11.3), Inches(3.9))
tf = tb.text_frame; tf.word_wrap = True
for i, line in enumerate(demo.split("\n")):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(1)
    r = p.add_run(); r.text = line if line else " "
    r.font.name = "Consolas"; r.font.size = Pt(12.5)
    r.font.color.rgb = MIST if not line.startswith("-") else RGBColor(0x8F, 0xE0, 0xE2)

# ============================================================ SLIDE 8 — LESSONS
s = prs.slides.add_slide(BLANK); bg(s, LIGHT)
header(s, "Reflection", "Lessons & Responsible AI", 8)
bullets(s, [
    "A tiny knowledge base + category boosts improved plans with no external model.",
    "Trace logging made a deterministic system feel far more trustworthy.",
    "Helpful AI moment: clean class design for the scheduler and task model.",
    "Flawed AI moment: proposed conflict detection on exact start-times only — "
    "missed overlaps; corrected with slot-based checks + a resolution pass.",
    "Known limit: categories are user-supplied, so a mislabel can misdirect a boost.",
], y=Inches(1.9), size=18, gap=13)
text(s, Inches(0.9), Inches(6.15), Inches(11.5), Inches(0.5),
     [("Full responsible-AI reflection in model_card.md", 12, TEAL, True)], space_after=0)

# ============================================================ SLIDE 9 — REPO
s = prs.slides.add_slide(BLANK); bg(s, LIGHT)
header(s, "Portfolio", "Repository & Files", 9)
files = [
    ("README.md", "project overview, setup, samples, testing"),
    ("model_card.md", "responsible-AI reflection"),
    ("pawpal_system.py", "model, scheduler, knowledge base, planner"),
    ("app.py  ·  main.py", "Streamlit UI & CLI demo"),
    ("tests/", "11 automated reliability checks"),
    ("diagrams/system_architecture.mmd", "architecture source"),
    ("ai_interactions.md", "agentic reasoning traces (stretch)"),
]
y = Inches(1.9)
for name, desc in files:
    rect(s, Inches(0.7), y, Inches(0.16), Inches(0.5), ORANGE)
    text(s, Inches(1.0), y - Inches(0.02), Inches(5.2), Inches(0.5),
         [(name, 15, NAVY, True)], space_after=0, anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(6.2), y - Inches(0.02), Inches(6.4), Inches(0.5),
         [(desc, 13, MUTE, False)], space_after=0, anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(0.62)
text(s, Inches(0.9), Inches(6.55), Inches(11.5), Inches(0.5),
     [("github.com/Giridhar555/ai110-module2show-pawpal-starter", 14, TEAL, True)],
     space_after=0)

# ============================================================ SLIDE 10 — THANKS
s = prs.slides.add_slide(BLANK); bg(s, NAVY)
rect(s, 0, Inches(3.35), SW, Inches(0.12), TEAL)
rect(s, 0, Inches(3.47), SW, Inches(0.06), ORANGE)
text(s, Inches(0), Inches(2.3), SW, Inches(1.1),
     [("Thank you  🐾", 52, WHITE, True)], align=PP_ALIGN.CENTER, space_after=0)
text(s, Inches(0), Inches(3.8), SW, Inches(0.8),
     [("Questions?", 24, MIST, False)], align=PP_ALIGN.CENTER, space_after=0)
text(s, Inches(0), Inches(4.6), SW, Inches(0.6),
     [("PawPal+  ·  Applied AI Pet-Care Planner  ·  Built by Giridhar", 15, TEAL, True)],
     align=PP_ALIGN.CENTER, space_after=0)

prs.save("PawPal_Plus_Presentation.pptx")
print("Saved PawPal_Plus_Presentation.pptx with", len(prs.slides._sldIdLst), "slides")
